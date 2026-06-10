"""
Validation chain for uploaded files.

Phase 1 — sequential, fail-fast (cheap checks):
  C1  size check
  C2  extension check
  C3  MIME type check
  C4  empty file check

Phase 2 — parallel (CPU/IO heavy):
  C5  corruption check  (PyMuPDF / python-pptx)
  C6  password check    (PyMuPDF)
  C7  duplicate check   (SHA-256 hash → Postgres via check_hash_fn callable)
  C8  virus scan        (ClamAV via pyclamd)
"""
from __future__ import annotations

import asyncio
import hashlib
import io
import zipfile
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path
from typing import Callable, Optional

import fitz  # PyMuPDF
import magic
import pyclamd
from loguru import logger
from pptx import Presentation

from app.config import settings
from app.upload_schemas import ValidatedFile


class ValidationError(Exception):
    def __init__(self, reason: str, stage: str, http_status: int = 400):
        super().__init__(reason)
        self.reason = reason
        self.stage = stage
        self.http_status = http_status


# ===========================================================================
# PHASE 1 — Sequential cheap checks
# ===========================================================================

def check_size(content_length: Optional[int], data: bytes) -> None:
    size = content_length if content_length is not None else len(data)
    max_bytes = settings.MAX_UPLOAD_SIZE_MB * 1024 * 1024
    if size > max_bytes:
        raise ValidationError(
            reason=f"File size {size} bytes exceeds limit of {settings.MAX_UPLOAD_SIZE_MB} MB",
            stage="size_check",
            http_status=413,
        )


def check_extension(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    allowed = {e.strip().lower() for e in settings.ALLOWED_EXTENSIONS.split(",")}
    if ext not in allowed:
        raise ValidationError(
            reason=f"Extension '{ext}' is not allowed. Allowed: {sorted(allowed)}",
            stage="extension_check",
            http_status=415,
        )
    return ext


def check_mime(data: bytes, ext: str) -> str:
    detected_mime = magic.from_buffer(data[:512], mime=True)
    mime_map = {}
    for pair in settings.MIME_MAP.split("|"):
        pair = pair.strip()
        if "=" in pair:
            k, v = pair.split("=", 1)
            mime_map[k.strip().lower()] = v.strip()
    expected_mime = mime_map.get(ext)
    if expected_mime and detected_mime != expected_mime:
        raise ValidationError(
            reason=f"MIME type '{detected_mime}' does not match extension '{ext}' (expected '{expected_mime}')",
            stage="mime_check",
        )
    return detected_mime


def _extract_text(data: bytes, ext: str) -> str:
    try:
        if ext == ".pdf":
            doc = fitz.open(stream=data, filetype="pdf")
            text = "".join(page.get_text() for page in doc)
            doc.close()
            return text
        elif ext == ".docx":
            import xml.etree.ElementTree as ET
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                with z.open("word/document.xml") as f:
                    tree = ET.parse(f)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            return " ".join(node.text or "" for node in tree.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"))
        elif ext in (".ppt", ".pptx"):
            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            return " ".join(parts)
    except Exception:
        return ""
    return ""


def check_empty(data: bytes, ext: str) -> None:
    # Rule 1: file must be greater than 0 bytes
    if len(data) == 0:
        raise ValidationError(
            reason="File is empty (0 bytes).",
            stage="empty_check",
        )

    # Rule 2: file must contain at least one non-whitespace character
    text = _extract_text(data, ext)
    if not text.strip():
        raise ValidationError(
            reason="File contains no readable text — only whitespace or no text layer.",
            stage="empty_check",
        )


# ===========================================================================
# PHASE 2 — Parallel deep checks
# ===========================================================================

def _check_corruption(data: bytes, ext: str) -> Optional[str]:
    try:
        if ext == ".pdf":
            doc = fitz.open(stream=data, filetype="pdf")
            doc.close()
        elif ext == ".docx":
            if not zipfile.is_zipfile(io.BytesIO(data)):
                return "File appears to be corrupted: not a valid DOCX (ZIP) file"
        elif ext in (".ppt", ".pptx"):
            Presentation(io.BytesIO(data))
        return None
    except Exception as e:
        return f"File appears to be corrupted: {e}"


def _check_password(data: bytes, ext: str) -> Optional[str]:
    if ext != ".pdf":
        return None
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        encrypted = doc.is_encrypted
        doc.close()
        return "File is password protected" if encrypted else None
    except Exception:
        return None


def _compute_sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _check_virus(data: bytes) -> Optional[str]:
    if not settings.CLAMAV_ENABLED:
        return None
    try:
        cd = pyclamd.ClamdNetworkSocket(host=settings.CLAMD_HOST, port=settings.CLAMD_PORT)
        result = cd.scan_stream(data)
        if result is None:
            return None
        threat = list(result.values())[0][1]
        return f"Virus/malware detected: {threat}"
    except pyclamd.ConnectionError:
        return None


# ===========================================================================
# Public entry point
# ===========================================================================

async def run_phase1(
    filename: str,
    data: bytes,
    check_hash_fn: Callable,
    log,
    content_length: Optional[int] = None,
) -> tuple[str, str, str, dict | None]:
    size = content_length if content_length is not None else len(data)

    try:
        check_size(content_length, data)
        log.info("  C1 size_check        PASS  ({:.2f} MB)", size / 1024 / 1024)
    except ValidationError:
        log.warning("  C1 size_check        FAIL")
        raise

    try:
        ext = check_extension(filename)
        log.info("  C2 extension_check   PASS  ({})", ext)
    except ValidationError as e:
        log.warning("  C2 extension_check   FAIL  ({})", e.reason)
        raise

    try:
        mime = check_mime(data, ext)
        log.info("  C3 mime_check        PASS  ({})", mime)
    except ValidationError as e:
        log.warning("  C3 mime_check        FAIL  ({})", e.reason)
        raise

    # C7 — duplicate check (requires SHA-256, placed here so Phase 2 only runs on unique files)
    sha256 = _compute_sha256(data)
    log.info("  SHA-256              {}", sha256)
    existing = await check_hash_fn(sha256)
    if existing:
        pdf_on_disk = Path(existing["pdf_path"]).exists()
        if pdf_on_disk:
            log.warning("  C7 duplicate_check   FAIL  (doc_id={} exists on disk)", existing["doc_id"])
            raise ValidationError(
                reason=f"Duplicate file — already uploaded as doc_id={existing['doc_id']}",
                stage="duplicate_check",
                http_status=409,
            )
        else:
            log.warning("  C7 duplicate_check   WARN  (in DB but missing from disk — re-upload allowed)")
            return ext, mime, sha256, existing
    else:
        log.info("  C7 duplicate_check   PASS  (new file)")

    return ext, mime, sha256, None


async def run_phase2(
    data: bytes,
    ext: str,
    log,
) -> None:
    log.info("  [C5 C6 C8] running in parallel ...")
    with ProcessPoolExecutor(max_workers=3) as pool:
        futures = {
            pool.submit(_check_corruption, data, ext): "C5 corruption_check",
            pool.submit(_check_password,   data, ext): "C6 password_check  ",
            pool.submit(_check_virus,      data):      "C8 virus_check     ",
        }
        errors: dict[str, str] = {}
        for future in as_completed(futures, timeout=settings.VALIDATION_PARALLEL_TIMEOUT):
            label = futures[future]
            try:
                result = future.result()
                if result is not None:
                    errors[label] = result
                    log.warning("  {} FAIL  ({})", label, result)
                else:
                    skipped = "(skipped — not PDF)" if "password" in label and ext != ".pdf" else ""
                    skipped = skipped or ("(skipped — ClamAV disabled)" if "virus" in label and not settings.CLAMAV_ENABLED else "")
                    log.info("  {} PASS  {}", label, skipped)
            except Exception as exc:
                errors[label] = str(exc)
                log.warning("  {} ERROR ({})", label, exc)

        if errors:
            stage, reason = next(iter(errors.items()))
            stage_key = stage.strip().split(" ")[-1]
            raise ValidationError(reason=reason, stage=stage_key)

    # C4 — empty check at end of Phase 2 (after expensive checks pass)
    try:
        check_empty(data, ext)
        log.info("  C4 empty_check       PASS  ({} bytes)", len(data))
    except ValidationError as e:
        log.warning("  C4 empty_check       FAIL  ({})", e.reason)
        raise


async def validate_upload(
    filename: str,
    data: bytes,
    check_hash_fn: Callable,
    log,
    content_length: Optional[int] = None,
) -> tuple[ValidatedFile, dict | None]:
    size_mb = len(data) / 1024 / 1024
    log.info("─── Validation started: {} ({:.2f} MB) ───", filename, size_mb)
    log.info("  Phase 1 — sequential checks")

    ext, mime, sha256, existing_doc = await run_phase1(filename, data, check_hash_fn, log, content_length)

    if existing_doc is not None:
        log.info("─── Validation complete: RE-UPLOAD (skipping Phase 2) ───")
        convertible = {e.strip().lower() for e in settings.CONVERTIBLE_EXTENSIONS.split(",")}
        return ValidatedFile(
            filename=filename, extension=ext, mime_type=mime,
            size_bytes=len(data), sha256_hash=sha256, content=data,
            needs_conversion=ext in convertible,
        ), existing_doc

    log.info("  Phase 2 — parallel deep checks")
    await run_phase2(data, ext, log)

    log.info("─── Validation complete: ALL CHECKS PASSED ───")

    convertible = {e.strip().lower() for e in settings.CONVERTIBLE_EXTENSIONS.split(",")}
    return ValidatedFile(
        filename=filename,
        extension=ext,
        mime_type=mime,
        size_bytes=len(data),
        sha256_hash=sha256,
        content=data,
        needs_conversion=ext in convertible,
    ), None
